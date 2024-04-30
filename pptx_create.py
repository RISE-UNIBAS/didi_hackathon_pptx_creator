import argparse
import os

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def create_pptx_with_images(args):
    """Create a PowerPoint presentation with two images per slide.
       Arguments are passed by the argparse object."""

    # Create a presentation object
    prs = Presentation()
    slide_width = prs.slide_width
    slide_width = slide_width - Inches(1)  # Subtract 1 inch to account for slide margins
    slide_height = prs.slide_height

    # Store the list of images in the second directory
    second_dir_list = os.listdir(args.second_dir)

    # Loop through each image in the directory
    image_idx = 0   # Index to keep track of the image in the second directory
    for image in os.listdir(args.first_dir):
        if args.verbose:
            print(f"Processing image: {image}")
        if image.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
            # Add a blank slide
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Calculate image path
            first_image = os.path.join(args.first_dir, image)
            second_image = os.path.join(args.second_dir, second_dir_list[image_idx])

            # Add and scale images
            for idx, image_path in enumerate([first_image, second_image]):
                img = Image.open(image_path)
                img_width, img_height = img.size

                # Calculate scaling to fit half the slide width while maintaining aspect ratio
                ratio = min(slide_width / 2 / img_width, slide_height / img_height)
                scaled_width = int(img_width * ratio)
                scaled_height = int(img_height * ratio)

                # Position images
                left = Inches(0.5) + idx * (slide_width / 2)  # Start first image at 0.5 inches, second at half slide width
                top = (slide_height - scaled_height) / 2  # Center vertically

                # Add image to slide
                slide.shapes.add_picture(image_path, left, top, width=scaled_width, height=scaled_height)
        else:
            if args.verbose:
                print(f"Skipping image: {image}")

        image_idx += 1

    # Save the presentation
    if args.verbose:
        print(f"Saving presentation...")
    prs.save(args.output_file_name)

    print(f"Presentation saved to {args.output_file_name}")


if __name__ == "__main__":
    # Create the parser
    parser = argparse.ArgumentParser(description='Create a PowerPoint presentation with two images per slide.')

    parser.add_argument("-v", "--verbose", action="store_true", help="Enable verbose mode")

    parser.add_argument('-o', '--output_file_name', type=str, help='Name of the output PowerPoint presentation',
                        default='output.pptx')

    parser.add_argument('first_dir', type=str, help='Directory containing the first set of images')
    parser.add_argument('second_dir', type=str, help='Directory containing the second set of images')

    # Parse the arguments
    p_args = parser.parse_args()

    # Perform the operation
    create_pptx_with_images(p_args)
