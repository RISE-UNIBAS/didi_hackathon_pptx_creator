import os
from pptx import Presentation
from pptx.util import Inches


def create_pptx_with_images(before_dir, after_dir, output_file_name):
    # Create a presentation object
    prs = Presentation()

    # Loop through each image in the directory
    image_idx = 0
    after_dir_list = os.listdir(after_dir)
    for image in os.listdir(before_dir):
        if image.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Adding a blank slide
            # Calculate image path
            image_path_before = os.path.join(before_dir, image)
            image_path_after = os.path.join(after_dir, after_dir_list[image_idx])
            # Add image to slide
            left = top = Inches(0)  # Positioning of the image

            pic = slide.shapes.add_picture(image_path_before, left, top, width=Inches(4), height=Inches(3))
            pic = slide.shapes.add_picture(image_path_after, Inches(5), top, width=Inches(4), height=Inches(3))

        image_idx += 1
    # Save the presentation
    prs.save(output_file_name)


# Example usage
before = 'images/before/'
after = 'images/after/'
output_pptx = 'output_presentation.pptx'
create_pptx_with_images(before, after, output_pptx)


